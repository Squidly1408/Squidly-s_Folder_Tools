# sections: folders, other

import tkinter as tk
from tkinter import filedialog, messagebox
from PyPDF2 import PdfReader, PdfWriter
from docx import Document
from pptx import Presentation
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
import os.path

# If modifying these SCOPES, delete the file token.json.
SCOPES = ["https://www.googleapis.com/auth/drive.readonly"]


class DocumentCombiner:
    def __init__(self, root):
        self.root = root
        self.root.title("Document Combiner")
        self.root.geometry("1000x600")
        self.root.configure(bg="#171717")

        # Custom Title Bar
        self.create_custom_title_bar()

        # Create Main Frame
        main_frame = tk.Frame(root, bg="#171717")
        main_frame.pack(fill="both", expand=True, padx=10, pady=10)

        # Number of Documents Input
        self.num_docs_var = tk.IntVar()
        num_docs_frame = tk.Frame(main_frame, bg="#171717")
        num_docs_frame.pack(pady=10)
        tk.Label(
            num_docs_frame,
            text="Number of Documents:",
            bg="#171717",
            fg="white",
            font=("Arial", 12),
        ).pack(side="left")
        tk.Entry(
            num_docs_frame, textvariable=self.num_docs_var, width=5, font=("Arial", 12)
        ).pack(side="left", padx=5)
        tk.Button(
            num_docs_frame,
            text="Load Documents",
            command=self.load_docs,
            font=("Arial", 12),
        ).pack(side="left", padx=5)

        # Frame for Document Pages Checkbuttons
        self.pages_frame = tk.Frame(main_frame, bg="#171717")
        self.pages_frame.pack(fill="both", expand=True, padx=10, pady=10)

        # Combine Button
        tk.Button(
            main_frame,
            text="Combine",
            command=self.combine_docs,
            font=("Arial", 12, "bold"),
            bg="#00796b",
            fg="white",
        ).pack(pady=10)

    def create_custom_title_bar(self):
        """Create a custom title bar with a close button."""
        title_bar = tk.Frame(self.root, bg="#00796b", relief="flat", height=30)
        title_bar.pack(fill="x", side="top")
        title_label = tk.Label(
            title_bar,
            text="Document Combiner",
            bg="#00796b",
            fg="white",
            font=("Arial", 12, "bold"),
        )
        title_label.pack(side="left", padx=10)
        close_button = tk.Button(
            title_bar,
            text="X",
            command=self.root.destroy,
            bg="#00796b",
            fg="white",
            font=("Arial", 10, "bold"),
            borderwidth=0,
            highlightthickness=0,
        )
        close_button.pack(side="right", padx=10)

    def load_docs(self):
        """Load documents and create checkbuttons for each page."""
        num_docs = self.num_docs_var.get()
        self.docs = []
        self.page_vars = []

        for i in range(num_docs):
            file_path = filedialog.askopenfilename(
                filetypes=[
                    ("PDF files", "*.pdf"),
                    ("Word files", "*.docx"),
                    ("Text files", "*.txt"),
                    ("PowerPoint files", "*.pptx"),
                    ("Google Docs", "*.gdoc"),
                    ("Google Slides", "*.gslides"),
                ]
            )
            if file_path:
                file_name = file_path.split("/")[-1]  # Extract the file name
                ext = file_name.split(".")[-1].lower()  # Get the file extension
                column = i  # Set the column for this document

                # Add a label at the top of the column
                tk.Label(
                    self.pages_frame,
                    text=file_name,
                    bg="#171717",
                    fg="white",
                    font=("Arial", 12, "bold"),
                ).grid(row=0, column=column, pady=5)

                if ext == "pdf":
                    reader = PdfReader(file_path)
                    self.docs.append((reader, "pdf"))
                    for page_num in range(len(reader.pages)):
                        var = tk.BooleanVar()
                        self.page_vars.append((var, "pdf", reader, page_num))
                        checkbutton = tk.Checkbutton(
                            self.pages_frame,
                            text=f"Page {page_num + 1}",
                            variable=var,
                            bg="#171717",
                            fg="white",
                            selectcolor="#00796b",
                            font=("Arial", 10),
                        )
                        checkbutton.grid(row=page_num + 1, column=column, sticky="w")
                elif ext == "docx":
                    doc = Document(file_path)
                    self.docs.append((doc, "docx"))
                    for page_num, paragraph in enumerate(doc.paragraphs):
                        var = tk.BooleanVar()
                        self.page_vars.append((var, "docx", doc, page_num))
                        checkbutton = tk.Checkbutton(
                            self.pages_frame,
                            text=f"Paragraph {page_num + 1}",
                            variable=var,
                            bg="#171717",
                            fg="white",
                            selectcolor="#00796b",
                            font=("Arial", 10),
                        )
                        checkbutton.grid(row=page_num + 1, column=column, sticky="w")
                elif ext == "txt":
                    with open(file_path, "r") as txt_file:
                        lines = txt_file.readlines()
                    self.docs.append((lines, "txt"))
                    for line_num, line in enumerate(lines):
                        var = tk.BooleanVar()
                        self.page_vars.append((var, "txt", lines, line_num))
                        checkbutton = tk.Checkbutton(
                            self.pages_frame,
                            text=f"Line {line_num + 1}",
                            variable=var,
                            bg="#171717",
                            fg="white",
                            selectcolor="#00796b",
                            font=("Arial", 10),
                        )
                        checkbutton.grid(row=line_num + 1, column=column, sticky="w")
                elif ext == "pptx":
                    ppt = Presentation(file_path)
                    self.docs.append((ppt, "pptx"))
                    for slide_num, slide in enumerate(ppt.slides):
                        var = tk.BooleanVar()
                        self.page_vars.append((var, "pptx", ppt, slide_num))
                        checkbutton = tk.Checkbutton(
                            self.pages_frame,
                            text=f"Slide {slide_num + 1}",
                            variable=var,
                            bg="#171717",
                            fg="white",
                            selectcolor="#00796b",
                            font=("Arial", 10),
                        )
                        checkbutton.grid(row=slide_num + 1, column=column, sticky="w")

    def combine_docs(self):
        """Combine selected pages from the loaded documents."""
        pdf_writer = PdfWriter()
        doc_combined = Document()
        txt_combined = []
        ppt_combined = Presentation()

        for var, doc_type, doc, index in self.page_vars:
            if var.get():
                if doc_type == "pdf":
                    pdf_writer.add_page(doc.pages[index])
                elif doc_type == "docx":
                    doc_combined.add_paragraph(doc.paragraphs[index].text)
                elif doc_type == "txt":
                    txt_combined.append(doc[index])
                elif doc_type == "pptx":
                    slide = ppt_combined.slides.add_slide(
                        doc.slides[index].slide_layout
                    )
                    for shape in doc.slides[index].shapes:
                        new_shape = slide.shapes.add_shape(
                            shape.auto_shape_type,
                            shape.left,
                            shape.top,
                            shape.width,
                            shape.height,
                        )
                        new_shape.text = shape.text

        output_path = filedialog.asksaveasfilename(
            defaultextension=".pdf",
            filetypes=[
                ("PDF files", "*.pdf"),
                ("Word files", "*.docx"),
                ("Text files", "*.txt"),
                ("PowerPoint files", "*.pptx"),
            ],
        )
        if output_path:
            ext = output_path.split(".")[-1].lower()
            if ext == "pdf":
                with open(output_path, "wb") as output_file:
                    pdf_writer.write(output_file)
            elif ext == "docx":
                doc_combined.save(output_path)
            elif ext == "txt":
                with open(output_path, "w") as output_file:
                    output_file.writelines(txt_combined)
            elif ext == "pptx":
                ppt_combined.save(output_path)
            messagebox.showinfo("Document Combiner", "Documents combined successfully!")


# Running the Document Combiner
if __name__ == "__main__":
    root = tk.Tk()
    app = DocumentCombiner(root)
    root.mainloop()
